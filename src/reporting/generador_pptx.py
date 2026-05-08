import ast
import calendar
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import ChartData
from pptx.enum.text import PP_ALIGN

# ── Paleta corporativa ──────────────────────────────────────────────────────
C_PRIMARY  = RGBColor(0,   82,  204)   # #0052CC
C_DARK     = RGBColor(30,  39,  46)
C_GRAY     = RGBColor(108, 117, 125)
C_LGRAY    = RGBColor(222, 226, 230)
C_LIGHT    = RGBColor(248, 249, 250)
C_ACCENT   = RGBColor(236, 243, 255)   # Fondo header suave
C_WHITE    = RGBColor(255, 255, 255)
C_SUCCESS  = RGBColor(40,  167, 69)
C_DANGER   = RGBColor(220, 53,  69)
C_WARN     = RGBColor(243, 156, 18)
C_EXT      = RGBColor(41,  128, 185)
C_INT      = RGBColor(230, 126, 34)
C_CAJA     = RGBColor(142, 68,  173)
C_OTRO     = RGBColor(127, 140, 141)

W = Inches(10)
H = Inches(7.5)

MESES_ES = {1:'Enero',2:'Febrero',3:'Marzo',4:'Abril',5:'Mayo',6:'Junio',
            7:'Julio',8:'Agosto',9:'Septiembre',10:'Octubre',11:'Noviembre',12:'Diciembre'}
DIAS_CORTO = ['L','M','X','J','V','S','D']


# ── Utilidades de zona ──────────────────────────────────────────────────────

def _color_zona(zona):
    z = str(zona).lower()
    if 'caja' in z:                     return C_CAJA
    if 'tienda' in z:                   return C_INT
    if 'exterior' in z or 'calle' in z: return C_EXT
    return C_OTRO

def _ordenar_zonas(zonas):
    def _peso(z):
        zl = str(z).lower()
        if 'exterior' in zl or 'calle' in zl: return 1
        if 'tienda' in zl:                    return 2
        if 'caja' in zl:                      return 3
        return 4
    return sorted(zonas, key=_peso)

def _zona_exterior(zonas):
    for z in _ordenar_zonas(zonas):
        zl = str(z).lower()
        if 'exterior' in zl or 'calle' in zl:
            return z
    return zonas[0] if zonas else None

def _zona_interior(zonas):
    for z in _ordenar_zonas(zonas):
        if 'tienda' in str(z).lower():
            return z
    return None

def _parsear_horas(val):
    try:
        r = ast.literal_eval(val) if isinstance(val, str) else val
        return r if isinstance(r, list) and len(r) == 24 else [0]*24
    except Exception:
        return [0]*24

def _fmt_num(v, decimales=0):
    if pd.isna(v): return '-'
    if decimales == 0:
        return f"{int(v):,}".replace(',', '.')
    return f"{v:,.{decimales}f}".replace(',', 'X').replace('.', ',').replace('X', '.')


# ── Gradiente para heatmap ───────────────────────────────────────────────────

def _interpolar_color(v, minv, maxv):
    if maxv <= minv or maxv == 0:
        return C_LIGHT
    t = max(0.0, min(1.0, (v - minv) / (maxv - minv)))
    r = int(248 - t * 248)
    g = int(249 - t * 167)
    b = int(250 - t * 46)
    return RGBColor(r, g, b)


# ── Primitivas de dibujo ─────────────────────────────────────────────────────

def _fondo(slide, color=C_WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def _rect(slide, x, y, w, h, color, sin_borde=True):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    if sin_borde:
        sh.line.fill.background()
    return sh

def _txt(slide, texto, x, y, w, h, size=12, bold=False, color=C_DARK,
         align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = str(texto)
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = align
    return tb

def _header(slide, titulo, subtitulo='', color_barra=C_PRIMARY):
    # Banda superior de color
    _rect(slide, 0, 0, W, Inches(0.34), color_barra)
    # Zona de título con fondo suave
    _rect(slide, 0, Inches(0.34), W, Inches(0.88), C_ACCENT)
    _txt(slide, titulo, Inches(0.4), Inches(0.38), Inches(9.2), Inches(0.58),
         size=22, bold=True, color=C_DARK)
    if subtitulo:
        _txt(slide, subtitulo, Inches(0.4), Inches(0.92), Inches(9.2), Inches(0.3),
             size=11, color=C_GRAY, italic=True)
    _rect(slide, 0, H - Inches(0.22), W, Inches(0.22), color_barra)

def _kpi_card(slide, x, y, w, h, etiqueta, valor, pie='', color=C_PRIMARY):
    _rect(slide, x, y, w, h, C_WHITE)
    # Borde izquierdo de color (más elegante que la barra superior)
    _rect(slide, x, y, Inches(0.07), h, color)
    offset = Inches(0.18)
    _txt(slide, etiqueta.upper(), x + offset, y + Inches(0.13),
         w - offset - Inches(0.08), Inches(0.26),
         size=9, bold=True, color=C_GRAY, align=PP_ALIGN.LEFT)
    _txt(slide, valor, x + offset, y + Inches(0.36),
         w - offset - Inches(0.08), Inches(0.7),
         size=26, bold=True, color=color, align=PP_ALIGN.LEFT, wrap=False)
    if pie:
        _txt(slide, pie, x + offset, y + h - Inches(0.3),
             w - offset - Inches(0.08), Inches(0.26),
             size=9, color=C_GRAY, align=PP_ALIGN.LEFT)

def _linea_seccion(slide, y, color=C_LGRAY):
    _rect(slide, Inches(0.4), y, Inches(9.2), Inches(0.02), color)


# ── Construcción de gráficos ─────────────────────────────────────────────────

def _chart_barras(slide, x, y, cx, cy, categorias, series_dict, titulo,
                  tipo=XL_CHART_TYPE.COLUMN_CLUSTERED, apilado=False):
    if not categorias or not series_dict:
        return
    cd = ChartData()
    cd.categories = [str(c) for c in categorias]
    for nombre, vals in series_dict.items():
        vals_clean = [float(v) if pd.notna(v) else 0.0 for v in vals]
        cd.add_series(nombre, vals_clean)

    if apilado:
        tipo = XL_CHART_TYPE.COLUMN_STACKED

    chart_sh = slide.shapes.add_chart(tipo, x, y, cx, cy, cd)
    chart = chart_sh.chart

    chart.has_legend = len(series_dict) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

    chart.chart_title.has_text_frame = True
    tf = chart.chart_title.text_frame
    tf.text = titulo
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = C_DARK

    nombres = list(series_dict.keys())
    for i, serie in enumerate(chart.series):
        c = _color_zona(nombres[i])
        serie.format.fill.solid()
        serie.format.fill.fore_color.rgb = c
        if not apilado:
            serie.format.line.fill.background()

    try:
        chart.value_axis.tick_labels.font.size = Pt(9)
        chart.category_axis.tick_labels.font.size = Pt(9)
    except Exception:
        pass

    return chart


def _chart_linea(slide, x, y, cx, cy, categorias, series_dict, titulo):
    return _chart_barras(slide, x, y, cx, cy, categorias, series_dict, titulo,
                         tipo=XL_CHART_TYPE.LINE)


# ── Heatmap calendario ───────────────────────────────────────────────────────

def _heatmap_mes(slide, x, y, cx, cy, df_mes, year, month):
    primer_dia = pd.Timestamp(year, month, 1)
    n_dias_mes = calendar.monthrange(year, month)[1]
    ultimo_dia = pd.Timestamp(year, month, n_dias_mes)

    df = df_mes.copy()
    df['fecha_d'] = pd.to_datetime(df['fecha']).dt.date
    por_dia = df.groupby('fecha_d')['total_visits'].sum()

    lunes_ini = primer_dia - pd.Timedelta(days=primer_dia.dayofweek)
    semanas = []
    cur = lunes_ini
    while cur <= ultimo_dia:
        semanas.append(cur)
        cur += pd.Timedelta(days=7)

    n_sem = len(semanas)
    n_cols = n_sem + 1
    n_rows = 8

    todos_vals = [por_dia.get(d, 0) for d in por_dia.index if d.month == month]
    maxv = max(todos_vals) if todos_vals else 1
    minv = min(v for v in todos_vals if v > 0) if todos_vals else 0

    tabla = slide.shapes.add_table(n_rows, n_cols, x, y, cx, cy).table

    ancho_etiq = Inches(0.45)
    ancho_col = int((cx - ancho_etiq) / n_sem)
    tabla.columns[0].width = ancho_etiq
    for j in range(1, n_cols):
        tabla.columns[j].width = ancho_col

    alto_cab = Inches(0.42)
    alto_fila = int((cy - alto_cab) / 7)
    tabla.rows[0].height = alto_cab
    for i in range(1, n_rows):
        tabla.rows[i].height = alto_fila

    _celda(tabla.cell(0, 0), '', C_DARK, C_WHITE, Pt(9), PP_ALIGN.CENTER)
    for j, lun in enumerate(semanas):
        _celda(tabla.cell(0, j+1),
               f"Sem {j+1}\n{lun.strftime('%d/%m')}",
               C_PRIMARY, C_WHITE, Pt(8), PP_ALIGN.CENTER, negrita=True)

    for i, dia in enumerate(DIAS_CORTO):
        _celda(tabla.cell(i+1, 0), dia, C_DARK, C_WHITE, Pt(10),
               PP_ALIGN.CENTER, negrita=True)
        for j, lun in enumerate(semanas):
            fecha_c = (lun + pd.Timedelta(days=i)).date()
            if fecha_c.month != month:
                _celda(tabla.cell(i+1, j+1), '', RGBColor(235, 235, 235), C_GRAY, Pt(8))
            else:
                v = int(por_dia.get(fecha_c, 0))
                if v > 0:
                    t = max(0.0, min(1.0, (v - minv) / (maxv - minv))) if maxv > minv else 0
                    bg = _interpolar_color(v, minv, maxv)
                    txt_c = C_WHITE if t > 0.52 else C_DARK
                else:
                    bg, txt_c = RGBColor(245, 245, 245), C_GRAY
                _celda(tabla.cell(i+1, j+1),
                       _fmt_num(v) if v > 0 else '–',
                       bg, txt_c, Pt(9), PP_ALIGN.CENTER)

def _celda(cell, texto, bg, fg, size, align=PP_ALIGN.LEFT, negrita=False):
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg
    cell.text = texto
    for para in cell.text_frame.paragraphs:
        para.alignment = align
        for run in para.runs:
            run.font.size = size
            run.font.bold = negrita
            run.font.color.rgb = fg


# ── Diapositivas individuales ────────────────────────────────────────────────

def _slide_portada(prs, blank, org_nombre, ubicaciones, periodo_str):
    s = prs.slides.add_slide(blank)
    _fondo(s, C_WHITE)

    # Sección superior azul (55% del slide)
    top_h = Inches(4.2)
    _rect(s, 0, 0, W, top_h, C_PRIMARY)

    # Etiqueta de categoría (texto pequeño en azul claro)
    _txt(s, 'INFORME DE ANÁLISIS DE AFLUENCIA',
         Inches(0.55), Inches(0.5), Inches(9), Inches(0.38),
         size=11, bold=True, color=RGBColor(140, 185, 255))

    # Nombre organización
    _txt(s, org_nombre or 'Informe Operativo',
         Inches(0.55), Inches(1.05), Inches(8.8), Inches(2.4),
         size=44, bold=True, color=C_WHITE)

    # Periodo en la zona azul
    _txt(s, periodo_str,
         Inches(0.55), Inches(3.5), Inches(7), Inches(0.45),
         size=16, color=RGBColor(190, 215, 255))

    # Zona blanca inferior
    _rect(s, Inches(0.55), top_h + Inches(0.45), Inches(1.2), Inches(0.05), C_PRIMARY)

    _txt(s, 'Emplazamientos analizados',
         Inches(0.55), top_h + Inches(0.6), Inches(5), Inches(0.28),
         size=9, bold=True, color=C_GRAY)

    ubi_txt = ' · '.join(ubicaciones[:4])
    if len(ubicaciones) > 4:
        ubi_txt += f' +{len(ubicaciones)-4} más'
    _txt(s, ubi_txt,
         Inches(0.55), top_h + Inches(0.9), Inches(9), Inches(0.42),
         size=14, bold=False, color=C_DARK)

    _txt(s, f'Generado el {pd.Timestamp("today").strftime("%d de %B de %Y")}',
         Inches(0.55), top_h + Inches(1.55), Inches(6), Inches(0.3),
         size=9, color=C_GRAY, italic=True)

    # Banda inferior oscura
    _rect(s, 0, H - Inches(0.38), W, Inches(0.38), C_DARK)
    _txt(s, 'Uso interno · Información confidencial · Aitanna Analytics',
         Inches(0.4), H - Inches(0.32), Inches(8), Inches(0.26),
         size=8, color=C_LGRAY, italic=True)


def _slide_vision_global(prs, blank, df, zonas, periodo_str):
    s = prs.slides.add_slide(blank)
    _fondo(s, C_WHITE)
    _header(s, f'Visión Global del Periodo',
            f'Agregado completo por zona · {periodo_str}')

    z_ext = _zona_exterior(zonas)
    z_int = _zona_interior(zonas)

    df_ext = df[df['Zona'] == z_ext] if z_ext else pd.DataFrame()
    df_int = df[df['Zona'] == z_int] if z_int else pd.DataFrame()

    total_v  = df_ext['total_visits'].sum()     if not df_ext.empty else df['total_visits'].sum()
    unicos_v = df_ext['unique_visitors'].mean()  if not df_ext.empty else df['unique_visitors'].mean()
    nuevos_v = df_ext['new_visitors'].sum()      if not df_ext.empty else df['new_visitors'].sum()
    estancia = (df_int['dwell_time'].mean()      if not df_int.empty
                else df['dwell_time'].mean()     if 'dwell_time' in df.columns else 0)

    card_w, card_h = Inches(2.15), Inches(1.35)
    gap = Inches(0.17)
    y_kpi = Inches(1.32)
    dt_g = f'{estancia:.1f} min'
    for i, (etiq, val, pie, col) in enumerate([
        ('Tráfico total',           _fmt_num(total_v),  z_ext or '',      C_EXT),
        ('Visitantes',              f'{unicos_v:.0f}',  'Media diaria',   C_INT),
        ('Nuevos visitantes',       _fmt_num(nuevos_v), 'Primera visita', C_SUCCESS),
        ('Estancia media',          dt_g,                z_int or '',     C_CAJA),
    ]):
        _kpi_card(s, Inches(0.4) + i*(card_w + gap), y_kpi, card_w, card_h,
                  etiq, val, pie, col)

    # Tabla resumen por zona
    kpi_cols = ['total_visits', 'unique_visitors', 'new_visitors', 'dwell_time']
    resumen = (df.groupby('Zona')
                 .agg({c: ('sum' if c != 'dwell_time' else 'mean')
                       for c in kpi_cols if c in df.columns})
                 .reset_index())
    resumen = resumen[resumen['Zona'].isin(zonas)]

    n_filas = len(resumen) + 1
    cabeceras = ['Zona', 'Tráfico total', 'Visitantes', 'Nuevos Visit.', 'Estancia Media']
    cols_tabla = [c for c in kpi_cols if c in resumen.columns]
    n_cols_t = 1 + len(cols_tabla)

    y_tabla = y_kpi + card_h + Inches(0.2)
    slide_tabla_resumen(s, Inches(0.4), y_tabla, Inches(9.2), Inches(0.42 * n_filas),
                        resumen, zonas, cabeceras[:n_cols_t], cols_tabla)

    # Gráfico: tendencia mensual
    df['YearMes'] = df['fecha'].dt.to_period('M').astype(str)
    meses_ord = sorted(df['YearMes'].unique())
    if len(meses_ord) > 1:
        y_chart = y_tabla + Inches(0.42 * n_filas) + Inches(0.2)
        h_chart = H - y_chart - Inches(0.38)
        series = {}
        for z in _ordenar_zonas(zonas):
            df_z = df[df['Zona'] == z].groupby('YearMes')['total_visits'].sum()
            series[z] = [int(df_z.get(m, 0)) for m in meses_ord]
        _chart_barras(s, Inches(0.4), y_chart, Inches(9.2), h_chart,
                      meses_ord, series, 'Tráfico total por zona y mes')


def slide_tabla_resumen(slide, x, y, cx, cy, resumen, zonas, cabeceras, cols_datos):
    n_filas = len(resumen) + 1
    n_cols = len(cabeceras)
    tbl = slide.shapes.add_table(n_filas, n_cols, x, y, cx, cy).table

    ancho_zona = Inches(2.2)
    ancho_col = int((cx - ancho_zona) / (n_cols - 1)) if n_cols > 1 else cx
    tbl.columns[0].width = ancho_zona
    for j in range(1, n_cols):
        tbl.columns[j].width = ancho_col

    alto_h = Inches(0.42)
    alto_f = Inches(0.38)
    tbl.rows[0].height = alto_h
    for i in range(1, n_filas):
        tbl.rows[i].height = alto_f

    for j, cab in enumerate(cabeceras):
        _celda(tbl.cell(0, j), cab, C_PRIMARY, C_WHITE, Pt(10), PP_ALIGN.CENTER, negrita=True)

    for i, row in resumen.iterrows():
        zona = row['Zona']
        bg = C_ACCENT if i % 2 == 0 else C_WHITE
        _celda(tbl.cell(i+1, 0), str(zona), bg, _color_zona(zona), Pt(10),
               PP_ALIGN.LEFT, negrita=True)
        for j, col in enumerate(cols_datos):
            if col not in resumen.columns:
                _celda(tbl.cell(i+1, j+1), '-', bg, C_DARK, Pt(10), PP_ALIGN.CENTER)
                continue
            v = row[col]
            if col == 'dwell_time':
                txt = f'{v:.1f} min' if pd.notna(v) else '-'
            else:
                txt = _fmt_num(v) if pd.notna(v) else '-'
            _celda(tbl.cell(i+1, j+1), txt, bg, C_DARK, Pt(10), PP_ALIGN.CENTER)
    return tbl


def _slide_kpis_mes(prs, blank, df_mes, zonas, nombre_mes, ubicacion=''):
    s = prs.slides.add_slide(blank)
    _fondo(s, C_WHITE)
    sub = ubicacion if ubicacion else ''
    _header(s, f'KPIs — {nombre_mes}', sub)

    z_ext = _zona_exterior(zonas)
    z_int = _zona_interior(zonas)
    df_e = df_mes[df_mes['Zona'] == z_ext] if z_ext else pd.DataFrame()
    df_i = df_mes[df_mes['Zona'] == z_int] if z_int else pd.DataFrame()

    tv   = df_e['total_visits'].sum()     if not df_e.empty else df_mes['total_visits'].sum()
    uv   = df_e['unique_visitors'].mean()  if not df_e.empty else df_mes['unique_visitors'].mean()
    nv   = df_e['new_visitors'].sum()      if not df_e.empty else df_mes['new_visitors'].sum()
    dt   = df_i['dwell_time'].mean()       if not df_i.empty else (df_mes['dwell_time'].mean() if 'dwell_time' in df_mes.columns else 0)
    uv7d = df_e['uv_7d'].mean()   if not df_e.empty and 'uv_7d'   in df_e.columns else None
    freq = df_e['freq_28d'].mean() if not df_e.empty and 'freq_28d' in df_e.columns else None

    card_w, card_h = Inches(2.15), Inches(1.35)
    gap = Inches(0.17)
    y_k = Inches(1.32)
    dt_fmt = f'{dt:.1f} min' if pd.notna(dt) else '-'
    for i, (etiq, val, pie, col) in enumerate([
        ('Tráfico total',           _fmt_num(tv),  f'{z_ext}' if z_ext else '',     C_EXT),
        ('Visitantes',              f'{uv:.0f}',   'Media diaria',                   C_INT),
        ('Nuevos visitantes',       _fmt_num(nv),  'Primera visita registrada',      C_SUCCESS),
        ('Estancia media',          dt_fmt,         f'{z_int}' if z_int else '',     C_CAJA),
    ]):
        _kpi_card(s, Inches(0.4) + i*(card_w+gap), y_k, card_w, card_h, etiq, val, pie, col)

    y_kpi2 = y_k + card_h + Inches(0.15)
    extras = []
    if uv7d is not None:
        extras.append(('Visitantes únicos 7 días (rolling)', _fmt_num(uv7d), C_INT))
    if freq is not None:
        extras.append(('Frecuencia de retorno (28d)', f'{freq:.2f}x', C_OTRO))
    for i, (etiq, val, color) in enumerate(extras[:2]):
        _kpi_card(s, Inches(0.4) + i*(Inches(4.4)+Inches(0.2)), y_kpi2,
                  Inches(4.4), Inches(1.0), etiq, val, color=color)

    y_chart = y_kpi2 + (Inches(1.15) if extras else Inches(0.15))
    h_chart = H - y_chart - Inches(0.38)

    df_mes = df_mes.copy()
    df_mes['fecha_str'] = df_mes['fecha'].dt.strftime('%d')
    dias_unicos = sorted(df_mes['fecha_str'].unique())
    cats = [f'{d}' for d in dias_unicos]

    zonas_ext_g = [z for z in _ordenar_zonas(zonas) if 'exterior' in z.lower() or 'calle' in z.lower()]
    zonas_int_g = [z for z in _ordenar_zonas(zonas) if z not in zonas_ext_g]

    if zonas_ext_g and zonas_int_g:
        w_chart = Inches(4.45)
        series_ext = {z: [int(df_mes[df_mes['Zona']==z].groupby('fecha_str')['total_visits'].sum().get(d, 0)) for d in dias_unicos] for z in zonas_ext_g}
        series_int = {z: [int(df_mes[df_mes['Zona']==z].groupby('fecha_str')['total_visits'].sum().get(d, 0)) for d in dias_unicos] for z in zonas_int_g}
        _chart_barras(s, Inches(0.4), y_chart, w_chart, h_chart, cats, series_ext, f'Exterior/Calle — {nombre_mes}')
        _chart_barras(s, Inches(5.15), y_chart, w_chart, h_chart, cats, series_int, f'Interior — {nombre_mes}')
    else:
        series = {z: [int(df_mes[df_mes['Zona']==z].groupby('fecha_str')['total_visits'].sum().get(d, 0)) for d in dias_unicos] for z in _ordenar_zonas(zonas)}
        _chart_barras(s, Inches(0.4), y_chart, Inches(9.2), h_chart, cats, series,
                      f'Tráfico diario por zona — {nombre_mes}')


def _slide_horario_mes(prs, blank, df_mes, zonas, nombre_mes):
    s = prs.slides.add_slide(blank)
    _fondo(s, C_WHITE)
    _header(s, f'Distribución Horaria — {nombre_mes}',
            'Tráfico acumulado por franja horaria · Horas de apertura (06:00 – 23:00)')

    horas_tot = [0] * 24
    df_mes = df_mes.copy()
    df_mes['hourly_array'] = df_mes['hourly_visits'].apply(_parsear_horas)
    for arr in df_mes['hourly_array']:
        horas_tot = [a + b for a, b in zip(horas_tot, arr)]

    rango = list(range(6, 24))
    cats  = [f'{h:02d}:00' for h in rango]
    vals  = [horas_tot[h] for h in rango]

    if sum(vals) == 0:
        return

    # Barras de columna: más legibles que línea cuando hay valores dispersos
    _chart_barras(s, Inches(0.4), Inches(1.32), Inches(9.2), Inches(4.85),
                  cats, {'Tráfico acumulado': vals},
                  f'Tráfico por franja horaria — {nombre_mes}',
                  tipo=XL_CHART_TYPE.COLUMN_CLUSTERED)

    h_max_idx = vals.index(max(vals))
    activos = [i for i, v in enumerate(vals) if v > 0]
    h_min_idx = min(activos, key=lambda i: vals[i]) if activos else 0
    _txt(s,
         f'Hora pico: {rango[h_max_idx]:02d}:00  ({_fmt_num(vals[h_max_idx])} pases acumulados)  ·  '
         f'Menor actividad: {rango[h_min_idx]:02d}:00',
         Inches(0.4), Inches(6.35), Inches(9.2), Inches(0.35),
         size=10, color=C_GRAY, italic=True)


def _slide_heatmap_mes(prs, blank, df_mes, zonas, nombre_mes, year, month):
    s = prs.slides.add_slide(blank)
    _fondo(s, C_WHITE)
    z_ref = _zona_exterior(zonas)
    zona_label = f'Zona: {z_ref}' if z_ref else ''
    _header(s, f'Mapa de Actividad Semanal — {nombre_mes}',
            f'Tráfico total (visitas) por día de la semana y semana del mes  ·  {zona_label}')

    df_ref = df_mes[df_mes['Zona'] == z_ref].copy() if z_ref else df_mes.copy()

    _heatmap_mes(s, Inches(0.4), Inches(1.32), Inches(9.2), Inches(4.25),
                 df_ref, year, month)

    legend_y = Inches(5.72)
    _txt(s, 'Intensidad:', Inches(0.4), legend_y, Inches(1.2), Inches(0.3),
         size=9, color=C_GRAY)
    colores_leyenda = [
        (RGBColor(242, 242, 242), 'Sin datos'),
        (_interpolar_color(0.1, 0, 1), 'Baja'),
        (_interpolar_color(0.4, 0, 1), 'Media'),
        (_interpolar_color(0.7, 0, 1), 'Alta'),
        (_interpolar_color(1.0, 0, 1), 'Máxima'),
    ]
    for k, (c, lbl) in enumerate(colores_leyenda):
        bx = Inches(1.7) + k * Inches(1.55)
        _rect(s, bx, legend_y + Inches(0.04), Inches(0.28), Inches(0.2), c)
        _txt(s, lbl, bx + Inches(0.34), legend_y, Inches(1.15), Inches(0.28),
             size=9, color=C_DARK)

    df_ref['diaw'] = pd.to_datetime(df_ref['fecha']).dt.dayofweek
    por_dia = df_ref.groupby('diaw')['total_visits'].sum()
    if not por_dia.empty:
        dia_max = por_dia.idxmax()
        dias_es = ['Lunes', 'Martes', 'Miércoles', 'Jueves', 'Viernes', 'Sábado', 'Domingo']
        _txt(s, f'Mayor afluencia: {dias_es[dia_max]} ({_fmt_num(int(por_dia[dia_max]))} visitas acumuladas en el mes)',
             Inches(0.4), Inches(6.32), Inches(9.2), Inches(0.34),
             size=10, color=C_GRAY, italic=True)


def _slide_ratio_atraccion(prs, blank, df_mes, zonas, nombre_mes):
    z_ext = _zona_exterior(zonas)
    z_int = _zona_interior(zonas)
    if not z_ext or not z_int:
        return

    df_e = df_mes[df_mes['Zona'] == z_ext].copy()
    df_i = df_mes[df_mes['Zona'] == z_int].copy()
    if df_e.empty or df_i.empty:
        return

    df_e['fecha_d'] = pd.to_datetime(df_e['fecha']).dt.date
    df_i['fecha_d'] = pd.to_datetime(df_i['fecha']).dt.date
    uv_e = df_e.groupby('fecha_d')['unique_visitors'].sum().reset_index(name='ext')
    uv_i = df_i.groupby('fecha_d')['unique_visitors'].sum().reset_index(name='int')
    merged = pd.merge(uv_e, uv_i, on='fecha_d')
    merged = merged[merged['ext'] > 0]
    if merged.empty:
        return
    merged['ratio'] = (merged['int'] / merged['ext'] * 100).round(1)

    s = prs.slides.add_slide(blank)
    _fondo(s, C_WHITE)
    _header(s, f'Ratio de Atracción — {nombre_mes}',
            f'{z_ext} → {z_int} · Porcentaje de visitantes exteriores que acceden al interior')

    ratio_medio = merged['ratio'].mean()
    ratio_max   = merged['ratio'].max()
    ratio_min   = merged['ratio'].min()

    color_ratio = C_SUCCESS if ratio_medio >= 50 else (C_WARN if ratio_medio >= 25 else C_DANGER)

    card_w, card_h = Inches(2.8), Inches(1.35)
    for i, (etiq, val, color) in enumerate([
        ('Ratio medio de atracción', f'{ratio_medio:.1f}%', color_ratio),
        ('Ratio máximo registrado',  f'{ratio_max:.1f}%',   C_SUCCESS),
        ('Ratio mínimo registrado',  f'{ratio_min:.1f}%',   C_DANGER),
    ]):
        _kpi_card(s, Inches(0.4) + i*(card_w + Inches(0.2)), Inches(1.32),
                  card_w, card_h, etiq, val, color=color)

    cats = [d.strftime('%d') for d in merged['fecha_d']]
    series = {f'Ratio {z_ext}→{z_int} (%)': merged['ratio'].tolist()}
    _chart_barras(s, Inches(0.4), Inches(2.88), Inches(9.2), Inches(4.0),
                  cats, series, f'Evolución diaria del ratio de atracción — {nombre_mes}',
                  tipo=XL_CHART_TYPE.LINE)


def _slide_kpis_zona_avanzados(prs, blank, df_mes, zonas, nombre_mes):
    cols_disponibles = [c for c in ['uv_7d', 'uv_28d', 'freq_7d', 'freq_28d']
                        if c in df_mes.columns]
    if not cols_disponibles:
        return

    s = prs.slides.add_slide(blank)
    _fondo(s, C_WHITE)
    _header(s, f'KPIs de Fidelización — {nombre_mes}',
            'Visitantes únicos acumulados y frecuencia de retorno por zona')

    etiquetas = {
        'uv_7d':   'Visitantes únicos (7d)',
        'uv_28d':  'Visitantes únicos (28d)',
        'freq_7d': 'Frecuencia retorno (7d)',
        'freq_28d':'Frecuencia retorno (28d)',
    }

    zonas_ord = _ordenar_zonas(zonas)
    n_cols_kpi = min(len(zonas_ord), 4)
    card_w = Inches(9.2 / n_cols_kpi) - Inches(0.15)
    card_h = Inches(1.1)
    gap = Inches(0.15)

    y_base = Inches(1.32)
    for row_idx, col_name in enumerate(cols_disponibles[:4]):
        for j, zona in enumerate(zonas_ord[:4]):
            df_z = df_mes[df_mes['Zona'] == zona]
            if df_z.empty:
                continue
            v = df_z[col_name].mean() if col_name.startswith('freq') else df_z[col_name].sum()
            val_str = f'{v:.2f}x' if col_name.startswith('freq') else _fmt_num(v)
            color = _color_zona(zona)
            _kpi_card(s,
                      Inches(0.4) + j*(card_w+gap),
                      y_base + row_idx*(card_h + Inches(0.1)),
                      card_w, card_h,
                      f'{etiquetas[col_name]} — {zona}',
                      val_str, color=color)


# ── Slide de conclusiones ────────────────────────────────────────────────────

def _slide_conclusiones(prs, blank, df, zonas, periodo_str, org_nombre=''):
    conclusiones = []

    z_ext = _zona_exterior(zonas)
    z_int = _zona_interior(zonas)
    df_e = df[df['Zona'] == z_ext].copy() if z_ext else df.copy()
    df_i = df[df['Zona'] == z_int].copy() if z_int else pd.DataFrame()

    df_e['fecha'] = pd.to_datetime(df_e['fecha'])

    # 1. Volumen global
    total = int(df_e['total_visits'].sum())
    avg_day = df_e.groupby(df_e['fecha'].dt.date)['total_visits'].sum().mean()
    conclusiones.append(('Volumen de afluencia',
        f"Durante el periodo se registraron {_fmt_num(total)} visitas totales, "
        f"con una media de {_fmt_num(int(avg_day))} visitas por día."))

    # 2. Distribución mensual
    df_e['YM'] = df_e['fecha'].dt.to_period('M')
    por_mes = df_e.groupby('YM')['total_visits'].sum()
    if len(por_mes) > 1:
        mes_max = por_mes.idxmax()
        mes_min = por_mes.idxmin()
        conclusiones.append(('Distribución mensual',
            f"{MESES_ES[mes_max.month]} fue el mes de mayor actividad ({_fmt_num(int(por_mes[mes_max]))} visitas). "
            f"{MESES_ES[mes_min.month]} registró la menor afluencia ({_fmt_num(int(por_mes[mes_min]))} visitas)."))

    # 3. Patrón semanal
    df_e['diaw'] = df_e['fecha'].dt.dayofweek
    por_dia = df_e.groupby('diaw')['total_visits'].mean()
    if not por_dia.empty:
        dias_es = ['Lunes', 'Martes', 'Miércoles', 'Jueves', 'Viernes', 'Sábado', 'Domingo']
        dia_max = por_dia.idxmax()
        dias_finde = [d for d in [5, 6] if d in por_dia.index]
        entre_sem = por_dia[[d for d in [0,1,2,3,4] if d in por_dia.index]].mean()
        fin_sem   = por_dia[dias_finde].mean() if dias_finde else np.nan
        if not np.isnan(fin_sem) and entre_sem > 0:
            ratio_fw = fin_sem / entre_sem
            patron = (f"El fin de semana genera {ratio_fw:.1f}x la afluencia media laboral." if ratio_fw > 1.1
                      else f"El perfil es marcadamente laboral: entre semana {1/ratio_fw:.1f}x más visitas que en fin de semana." if ratio_fw < 0.9
                      else "La afluencia se distribuye de manera homogénea a lo largo de toda la semana.")
        else:
            patron = ''
        conclusiones.append(('Patrón semanal',
            f"El {dias_es[dia_max]} concentra la mayor afluencia media del periodo "
            f"({_fmt_num(int(por_dia[dia_max]))} visitas/día). {patron}"))

    # 4. Ratio de atracción
    if z_ext and z_int and not df_i.empty:
        df_e2 = df_e.copy()
        df_i2 = df_i.copy()
        df_i2['fecha'] = pd.to_datetime(df_i2['fecha'])
        df_e2['fd'] = df_e2['fecha'].dt.date
        df_i2['fd'] = df_i2['fecha'].dt.date
        uv_e = df_e2.groupby('fd')['unique_visitors'].sum()
        uv_i = df_i2.groupby('fd')['unique_visitors'].sum()
        merged = pd.merge(uv_e, uv_i, left_index=True, right_index=True, suffixes=('_e','_i'))
        merged = merged[merged['unique_visitors_e'] > 0]
        if not merged.empty:
            ratio = (merged['unique_visitors_i'] / merged['unique_visitors_e'] * 100).mean()
            nivel = 'alta' if ratio >= 50 else ('moderada' if ratio >= 25 else 'baja')
            accion = ('Mantener la propuesta de valor en escaparate y acciones de fidelización.' if ratio >= 50
                      else 'Reforzar la señalización exterior y activar acciones de captación en puerta.' if ratio >= 25
                      else 'Se recomienda revisar la visibilidad del establecimiento y activar campañas de captación.')
            conclusiones.append(('Ratio de atracción exterior → interior',
                f"De cada 100 personas que pasan por la zona exterior, {int(ratio)} acceden al interior "
                f"(ratio {ratio:.1f}%, conversión {nivel}). {accion}"))

    # 5. Captación vs. fidelización
    if 'new_visitors' in df_e.columns and 'unique_visitors' in df_e.columns:
        nv = int(df_e['new_visitors'].sum())
        uv = df_e['unique_visitors'].sum()
        if uv > 0 and nv > 0:
            pct = nv / uv * 100
            tipo = ('alta captación de nuevos clientes' if pct > 60
                    else 'equilibrio entre captación y fidelización' if pct > 40
                    else 'base de clientes fidelizada con alta recurrencia')
            conclusiones.append(('Captación vs. fidelización',
                f"El {pct:.1f}% de los visitantes únicos son nuevos ({_fmt_num(nv)} primeras visitas). "
                f"Perfil de {tipo}."))

    # 6. Estancia media (comportamiento interior)
    if not df_i.empty and 'dwell_time' in df_i.columns:
        df_i2 = df_i.copy()
        dt = df_i2['dwell_time'].mean()
        if pd.notna(dt) and dt > 0:
            nivel_dt = 'alta implicación con el espacio' if dt > 15 else ('exploración moderada' if dt > 8 else 'visita transaccional de corta duración')
            recom = ('El tiempo de permanencia es favorable para la conversión.' if dt >= 8
                     else 'Oportunidad de mejora en la experiencia de compra para prolongar la estancia.')
            conclusiones.append(('Comportamiento en el interior',
                f"La estancia media en {z_int} es de {dt:.1f} min, indicador de {nivel_dt}. {recom}"))

    # 7. Tendencia del periodo
    df_trend = df_e.sort_values('fecha')
    n = len(df_trend)
    if n >= 30:
        tercio = n // 3
        v1 = df_trend.iloc[:tercio]['total_visits'].mean()
        v3 = df_trend.iloc[-tercio:]['total_visits'].mean()
        if v1 > 0:
            var_pct = (v3 - v1) / v1 * 100
            if var_pct > 10:
                tend_txt = f"tendencia ascendente (+{var_pct:.1f}%): la afluencia ha mejorado hacia el cierre del periodo."
            elif var_pct < -10:
                tend_txt = f"tendencia descendente ({var_pct:.1f}%): la afluencia ha disminuido en la segunda parte del periodo."
            else:
                tend_txt = f"comportamiento estable (variación del {var_pct:+.1f}% entre el inicio y el cierre del periodo)."
            conclusiones.append(('Evolución temporal', f"Se detecta una {tend_txt}"))

    # ── Render ─────────────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _fondo(s, C_WHITE)
    _header(s, 'Conclusiones y Hallazgos Analíticos',
            f'Síntesis automática del periodo · {periodo_str}', color_barra=C_DARK)

    n_c = min(len(conclusiones), 7)
    espacio_disp = H - Inches(1.38) - Inches(0.38)
    h_bloque = min(espacio_disp / max(n_c, 1), Inches(0.88))
    COLS = [C_EXT, C_INT, C_PRIMARY, C_SUCCESS, C_WARN, C_CAJA, C_DANGER]

    y = Inches(1.42)
    for i, (titulo_c, texto_c) in enumerate(conclusiones[:n_c]):
        c = COLS[i % len(COLS)]
        _rect(s, Inches(0.4), y + Inches(0.05), Inches(0.06), h_bloque - Inches(0.1), c)
        _txt(s, titulo_c.upper(), Inches(0.6), y + Inches(0.04),
             Inches(8.9), Inches(0.26),
             size=9, bold=True, color=c)
        _txt(s, texto_c, Inches(0.6), y + Inches(0.28),
             Inches(8.9), h_bloque - Inches(0.32),
             size=10, color=C_DARK)
        y += h_bloque

    _txt(s,
         'Análisis generado automáticamente a partir de los datos históricos del sistema · Para uso interno',
         Inches(0.4), H - Inches(0.55), Inches(9.2), Inches(0.28),
         size=8, color=C_GRAY, italic=True, align=PP_ALIGN.CENTER)


# ── Función principal ─────────────────────────────────────────────────────────

def generar_reporte_pptx(df, stream, start_date, end_date, org_nombre=''):
    prs = Presentation()
    blank = prs.slide_layouts[6]

    df = df.copy()
    df['fecha'] = pd.to_datetime(df['fecha'])
    df['YearMes'] = df['fecha'].dt.to_period('M')
    df['Año'] = df['fecha'].dt.year
    df['Mes'] = df['fecha'].dt.month

    df = df[df['Zona'] != 'SinNombre']
    df = df[~df['Zona'].str.contains('Extra|End|exit', case=False, na=False)]

    zonas_globales = _ordenar_zonas(df['Zona'].unique().tolist())
    ubicaciones    = df['Ubicación'].unique().tolist()
    periodo_str    = f'{start_date.strftime("%d/%m/%Y")} — {end_date.strftime("%d/%m/%Y")}'

    # ── PORTADA ───────────────────────────────────────────────────────────────
    titulo_portada = (org_nombre or
                      (ubicaciones[0] if len(ubicaciones) == 1 else f'{len(ubicaciones)} Emplazamientos'))
    _slide_portada(prs, blank, titulo_portada, ubicaciones, periodo_str)

    # ── VISIÓN GLOBAL ─────────────────────────────────────────────────────────
    _slide_vision_global(prs, blank, df, zonas_globales, periodo_str)

    # ── DIAPOSITIVAS MENSUALES ────────────────────────────────────────────────
    meses_presentes = sorted(df[['Año', 'Mes']].drop_duplicates().itertuples(index=False),
                             key=lambda r: (r.Año, r.Mes))

    for registro in meses_presentes:
        year, mes = registro.Año, registro.Mes
        nombre_mes = f'{MESES_ES[mes]} {year}'
        df_mes = df[(df['Año'] == year) & (df['Mes'] == mes)].copy()
        if df_mes.empty:
            continue

        for ubi in df_mes['Ubicación'].unique():
            df_ubi = df_mes[df_mes['Ubicación'] == ubi].copy()
            zonas_ubi = _ordenar_zonas(df_ubi['Zona'].unique().tolist())
            sub_ubi = ubi if len(ubicaciones) > 1 else ''

            _slide_kpis_mes(prs, blank, df_ubi, zonas_ubi, nombre_mes, sub_ubi)

            if 'hourly_visits' in df_ubi.columns:
                _slide_horario_mes(prs, blank, df_ubi, zonas_ubi, nombre_mes)

            _slide_heatmap_mes(prs, blank, df_ubi, zonas_ubi, nombre_mes, year, mes)

            if _zona_exterior(zonas_ubi) and _zona_interior(zonas_ubi):
                _slide_ratio_atraccion(prs, blank, df_ubi, zonas_ubi, nombre_mes)

            _slide_kpis_zona_avanzados(prs, blank, df_ubi, zonas_ubi, nombre_mes)

    # ── CONCLUSIONES ──────────────────────────────────────────────────────────
    _slide_conclusiones(prs, blank, df, zonas_globales, periodo_str, org_nombre)

    prs.save(stream)
